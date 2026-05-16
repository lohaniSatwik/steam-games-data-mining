#!/usr/bin/env python3
"""
Generate steam_games_presentation.pptx  (12 slides, ~10 min)
IE500 Data Mining | Team 9 – Brewed Clusters
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

FIGURES = os.path.join('latex_report', 'figures')

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(31,  56,  100)
STEEL_BLUE = RGBColor(70,  130, 180)
LT_BLUE    = RGBColor(219, 234, 254)
WHITE      = RGBColor(255, 255, 255)
DARK       = RGBColor(30,  30,  30)
GREY       = RGBColor(110, 110, 110)
GREEN      = RGBColor(21,  128, 61)
ORANGE     = RGBColor(194, 100, 20)
LT_GREEN   = RGBColor(220, 252, 231)
LT_ORANGE  = RGBColor(254, 243, 199)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── Primitive helpers ─────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)


def box(slide, l, t, w, h, fill=None, no_line=True):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if no_line:
        s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h, size=16, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return shape


def header_bar(slide, title, subtitle=None):
    box(slide, Inches(0), Inches(0), SW, Inches(1.25), fill=DARK_BLUE)
    box(slide, Inches(0), Inches(1.25), SW, Inches(0.055), fill=STEEL_BLUE)
    tb(slide, title,
       Inches(0.35), Inches(0.08), Inches(12.6), Inches(0.85),
       size=28, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.35), Inches(0.88), Inches(12.6), Inches(0.38),
           size=13, color=LT_BLUE)


def bullets(slide, items, l, t, w, h, size=16):
    """items: list of str.  Prefix '  ' for sub-bullet."""
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        sub = item.startswith('  ')
        text = item.strip()
        bullet = '–  ' if sub else '•  '
        r = p.add_run()
        r.text = bullet + text
        r.font.size = Pt(size - 2 if sub else size)
        r.font.color.rgb = GREY if sub else DARK
        # paragraph spacing after
        pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spc    = etree.SubElement(spcAft, qn('a:spcPts'))
        spc.set('val', '300' if not sub else '150')
    return shape


def pic(slide, fname, l, t, w):
    path = os.path.join(FIGURES, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, width=w)
    else:
        box(slide, l, t, w, Inches(3), fill=LT_BLUE)
        tb(slide, f'[{fname}]', l, t + Inches(1.2), w, Inches(0.6),
           size=11, color=GREY, align=PP_ALIGN.CENTER)


def footer(slide, n, total=12):
    box(slide, Inches(0), SH - Inches(0.32), SW, Inches(0.32), fill=DARK_BLUE)
    tb(slide, f'IE500 Data Mining | Team 9 – Brewed Clusters | {n}/{total}',
       Inches(0.4), SH - Inches(0.30), Inches(12.5), Inches(0.28),
       size=10, color=LT_BLUE, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
box(s, Inches(0), Inches(0),    SW, Inches(0.2),  fill=DARK_BLUE)
box(s, Inches(0), Inches(7.3),  SW, Inches(0.2),  fill=DARK_BLUE)
box(s, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.05), fill=STEEL_BLUE)

tb(s, 'Predicting Game Success on the Steam Platform',
   Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5),
   size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
tb(s, 'Multiclass Classification of Steam Game Reception',
   Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.65),
   size=20, color=GREY, align=PP_ALIGN.CENTER)
tb(s, 'Team 9 – Brewed Clusters',
   Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.55),
   size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
tb(s, 'Satwik Kumar Lohani  ·  Arpitha Shivaprasad Kori  ·  Zubair Ashfaq  ·  '
      'Beyza Ünlü  ·  Muhammad Sameer Siddiqui  ·  Esha Raheel',
   Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.9),
   size=14, color=GREY, align=PP_ALIGN.CENTER)
tb(s, 'IE500 Data Mining  |  University of Mannheim  |  Spring 2026',
   Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.5),
   size=13, color=GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Problem Statement', 'Why is predicting Steam game reception difficult?')
footer(s, 2)

bullets(s, [
    'Steam hosts 50,000+ active titles — developers need early signals of reception before release',
    'Task: predict Good / Mixed / Bad reception using only pre-release metadata',
    '  Good: review ratio ≥ 0.75   |   Mixed: 0.50 – 0.74   |   Bad: < 0.50',
    'Primary metric: Macro F1 — penalises poor minority-class performance equally',
    'Key challenge: 3 imbalanced classes (Good 63 %, Mixed 28 %, Bad 8 %)',
    '  Good and Mixed games share nearly identical metadata; only review ratio differs',
    '  Current features cannot encode execution quality, post-launch patches, or community sentiment',
    'Leakage constraint: must exclude post-release signals (playtime, Metacritic, Recommendations)',
], Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.6), size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset & Target
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Dataset & Target Construction', 'Steam Games Dataset — Kaggle (FronkonGames, 2024)')
footer(s, 3)

# Left column: pipeline
bullets(s, [
    'Raw CSV: 122,611 rows × 39 columns',
    '  Price, genres, tags, categories, platform flags, achievements, DLC count',
    'Review-stability filter: keep games with ≥ 10 total reviews',
    '  56,655 rows retained',
    'Stratified 80 / 20 train–test split',
    '  Train: 45,324 games  |  Test: 11,331 games',
    'Safe feature set: 147 columns (no post-release leakage)',
], Inches(0.5), Inches(1.45), Inches(6.2), Inches(4.5), size=16)

# Right column: class distribution table
box(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(4.2), fill=LT_BLUE)
tb(s, 'Class Distribution (Train)',
   Inches(7.2), Inches(1.55), Inches(5.4), Inches(0.5),
   size=16, bold=True, color=DARK_BLUE)
rows = [
    ('Good',  'ratio ≥ 0.75', '28,671', '63.3 %', LT_GREEN),
    ('Mixed', '0.50 – 0.74',  '12,834', '28.3 %', LT_ORANGE),
    ('Bad',   'ratio < 0.50', ' 3,819',  '8.4 %', RGBColor(254,226,226)),
]
for i, (cls, rule, cnt, pct, bg) in enumerate(rows):
    y = Inches(2.15) + i * Inches(0.9)
    box(s, Inches(7.05), y, Inches(5.7), Inches(0.82), fill=bg)
    tb(s, f'{cls}', Inches(7.15), y + Inches(0.1), Inches(1.3), Inches(0.6),
       size=17, bold=True, color=DARK)
    tb(s, rule, Inches(8.5), y + Inches(0.1), Inches(2.0), Inches(0.6),
       size=14, color=GREY)
    tb(s, f'{cnt}  ({pct})', Inches(10.5), y + Inches(0.1), Inches(2.0), Inches(0.6),
       size=14, bold=True, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Preprocessing
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Preprocessing Pipeline', '9 steps applied identically to train and test splits')
footer(s, 4)

bullets(s, [
    'Duplicate & validity checks — exact duplicates removed; negative prices/counts corrected',
    'Row-alignment fix — misaligned CSV rows corrected before any downstream processing',
    'Numeric conversion — price, positive/negative review counts, achievements cast to float',
    'Weak-column removal — AppIDs, URLs, review text, support links dropped',
    'Review-stability filter — games with < 10 reviews excluded (→ 56,655 rows)',
    'Multi-value parsing — genres, categories, tags binarised via MultiLabelBinarizer',
    '  Top 50 tags by training-set frequency retained; 58 category flags; 28 genre flags',
    'Derived features — log_price, release era buckets (pre-2015 / 2015-2019 / 2020+), n_languages',
    'Leakage-aware split — avg_playtime, Recommendations, Metacritic excluded from safe feature set',
    'Scaling — StandardScaler fitted on train only; applied to all continuous features',
], Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.5), size=16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Exploratory Data Analysis', 'Four modelling-relevant patterns in the safe training set')
footer(s, 5)

# Four highlighted finding boxes
findings = [
    (DARK_BLUE, WHITE, 'Class Imbalance',
     'Good (63.3%) dominates.\nBad (8.4%) is the minority — limited training signal.'),
    (STEEL_BLUE, WHITE, 'Sparse Features',
     'Most binary indicators are True for < 20% of games.\nIndividual features have limited separating power.'),
    (RGBColor(21,128,61), WHITE, 'Weak Numeric Correlations',
     'Pearson correlations between safe numerics and\nGood-vs-rest are all below |0.15|.\nNon-linear combinations are required.'),
    (RGBColor(146,64,14), WHITE, 'Descriptor Patterns',
     'cat_Steam Cloud, cat_Full controller support,\nera_2020+ → higher good-reception rates.\nMonetisation-heavy categories trend lower.'),
]

for i, (bg, fg, title, body) in enumerate(findings):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.5)
    y = Inches(1.5  + row * 2.85)
    box(s, x, y, Inches(6.2), Inches(2.65), fill=bg)
    tb(s, title, x + Inches(0.2), y + Inches(0.15),
       Inches(5.8), Inches(0.5), size=17, bold=True, color=fg)
    tb(s, body, x + Inches(0.2), y + Inches(0.65),
       Inches(5.8), Inches(1.85), size=14, color=fg, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Modelling Setup
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Modelling Setup', '7 models trained under identical nested cross-validation')
footer(s, 6)

bullets(s, [
    'Nested cross-validation: 5 outer folds (unbiased performance) + 3 inner folds (GridSearchCV tuning)',
    'Primary metric throughout: Macro F1 — weights all three classes equally',
    'Class imbalance handled inside each model pipeline:',
    '  class_weight="balanced" for RF, XGBoost, SVM, LR baseline',
    '  SMOTE and Random Over-Sampling tested separately on LR',
    'Final model: GridSearchCV on full training set → evaluate on test set once',
], Inches(0.5), Inches(1.45), Inches(12.3), Inches(3.0), size=17)

# Models list
box(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), fill=LT_BLUE)
tb(s, 'Models evaluated:', Inches(0.7), Inches(4.3), Inches(12.0), Inches(0.5),
   size=16, bold=True, color=DARK_BLUE)
model_cols = [
    ['Random Forest (RF)', 'XGBoost', 'Linear SVM'],
    ['Logistic Regression (LR)', 'LR + SMOTE', 'LR + Random Over-Sampling'],
    ['K-Nearest Neighbours (KNN)', '', 'Majority-class baseline'],
]
for ci, col in enumerate(model_cols):
    for ri, m in enumerate(col):
        if m:
            tb(s, '•  ' + m,
               Inches(0.7 + ci * 4.2), Inches(4.85) + ri * Inches(0.58),
               Inches(4.1), Inches(0.55), size=15, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — All Model Results
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Model Comparison Results', 'Random Forest achieves the best Test Macro F1 = 0.5100')
footer(s, 7)

pic(s, 'fig03_model_comparison.png', Inches(0.4), Inches(1.4), Inches(7.5))

# Summary table on right
box(s, Inches(8.1), Inches(1.4), Inches(5.0), Inches(5.5), fill=LT_BLUE)
tb(s, 'Test Macro F1 Summary',
   Inches(8.2), Inches(1.5), Inches(4.8), Inches(0.45),
   size=14, bold=True, color=DARK_BLUE)
table_rows = [
    ('Random Forest',    '0.5100', GREEN),
    ('XGBoost',          '0.4799', DARK),
    ('Linear SVM',       '0.4693', DARK),
    ('LR + SMOTE',       '0.4648', DARK),
    ('KNN',              '0.4516', DARK),
    ('LR baseline',      '0.4355', DARK),
    ('LR + ROS',         '0.4352', DARK),
    ('Majority class',   '0.2583', ORANGE),
]
for i, (name, score, color) in enumerate(table_rows):
    y = Inches(2.05) + i * Inches(0.52)
    if i == 0:
        box(s, Inches(8.15), y - Inches(0.03), Inches(4.9), Inches(0.5), fill=LT_GREEN)
    tb(s, name,  Inches(8.25), y, Inches(3.2), Inches(0.48), size=13, color=color)
    tb(s, score, Inches(11.5), y, Inches(1.4), Inches(0.48), size=13, bold=(i==0), color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Best Model Deep-Dive
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Best Model: Random Forest', 'Test Macro F1 = 0.5100  |  n_estimators=200, max_depth=None, min_samples_leaf=4')
footer(s, 8)

pic(s, 'fig05_rf_cm.png', Inches(0.4), Inches(1.4), Inches(5.8))

bullets(s, [
    'Good F1 = 0.77  |  Mixed F1 = 0.43  |  Bad F1 = 0.33',
    'Tree ensembles outperform linear models because feature interactions are non-linear',
    '  e.g. tag_2D × era_2020+ is disproportionately Good — LR cannot model this directly',
    'RF beats XGBoost on balanced Macro F1:',
    '  XGBoost recovers higher Bad recall (0.56) but lower Good F1 → net lower macro',
    'class_weight="balanced" with max_features="sqrt" gives best bias–variance trade-off',
    'Consistent performance: Nested CV F1 = 0.5093 ± 0.0073',
], Inches(6.6), Inches(1.45), Inches(6.4), Inches(5.5), size=15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Error Analysis
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Error Analysis', 'Why do Mixed (F1=0.43) and Bad (F1=0.33) underperform?')
footer(s, 9)

pic(s, 'fig15b_confidence.png', Inches(0.4), Inches(1.4), Inches(6.5))

# Error flow table on right
box(s, Inches(7.1), Inches(1.4), Inches(6.0), Inches(2.6), fill=LT_BLUE)
tb(s, 'Dominant error patterns',
   Inches(7.3), Inches(1.5), Inches(5.6), Inches(0.45),
   size=14, bold=True, color=DARK_BLUE)
err_rows = [
    ('Mixed → Good', '1,360', '42.4 % of Mixed', ORANGE),
    ('Mixed → Bad',  '  505', '16.6 % of Mixed', DARK),
    ('Bad → Mixed',  '  358', '37.5 % of Bad',   DARK),
    ('Bad → Good',   '  229', '24.0 % of Bad',   DARK),
]
for i, (flow, cnt, pct, color) in enumerate(err_rows):
    y = Inches(2.05) + i * Inches(0.52)
    if i == 0:
        box(s, Inches(7.15), y - Inches(0.03), Inches(5.9), Inches(0.5), fill=LT_ORANGE)
    tb(s, flow, Inches(7.2), y, Inches(2.6), Inches(0.48), size=13, color=color, bold=(i==0))
    tb(s, cnt,  Inches(9.8), y, Inches(0.8), Inches(0.48), size=13, color=color)
    tb(s, pct,  Inches(10.6),y, Inches(2.3), Inches(0.48), size=13, color=color)

bullets(s, [
    'Correct predictions: median confidence = 0.570',
    'Wrong predictions:   median confidence = 0.460',
    '320 high-confidence errors (prob ≥ 0.70) — 75.9 % are Mixed→Good',
    'Root cause: Good and Mixed games are genuinely similar in feature space',
    '  Only review ratio separates them — no pre-release metadata encodes this gap',
], Inches(7.1), Inches(4.15), Inches(6.0), Inches(3.0), size=15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Feature Importance & Separation
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Feature Importance & Good–Mixed Separation',
           'Even the strongest separating signal (+0.172) is modest')
footer(s, 10)

pic(s, 'fig15c_feat_separation.png', Inches(0.4), Inches(1.4), Inches(7.8))

bullets(s, [
    'Top Good − Mixed feature differences:',
    '  tag_2D:                  +0.172',
    '  era_2020+:               +0.159',
    '  cat_Steam Cloud:         +0.145',
    '  era_2015-2019:           −0.144',
    '  tag_Singleplayer:        +0.128',
    '  cat_Full controller:     +0.122',
    'Implication: modern, polished, singleplayer 2D games trend Good',
    'But no single feature cleanly separates the classes — differences are modest',
    'LR symbolic analysis confirms: era_2020+, tag_2D, cat_Steam Cloud are top positive LR coefficients',
], Inches(8.4), Inches(1.45), Inches(4.7), Inches(5.5), size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Feature Engineering
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Feature Engineering Experiment',
           '8 derived features added (147 → 155) — no meaningful improvement')
footer(s, 11)

pic(s, 'fig17c_feat_importance.png', Inches(0.4), Inches(1.4), Inches(6.8))

# Features list
fe_items = [
    ('n_tags',           'count',        'More tags = more discoverable'),
    ('n_categories',     'count',        'Release polish proxy'),
    ('n_genres',         'count',        'Genre breadth'),
    ('polish_index',     'composite',    'Steam Cloud + controller + achievements + …'),
    ('recent_2d',        'binary',       'tag_2D × era_2020+'),
    ('price_x_recent',   'interaction',  'log_price × era_2020+'),
    ('solo_only',        'binary',       'Singleplayer and not Multiplayer'),
    ('platform_breadth', 'count',        'Windows + Mac + Linux'),
]
box(s, Inches(7.4), Inches(1.4), Inches(5.7), Inches(3.8), fill=LT_BLUE)
tb(s, 'New features', Inches(7.55), Inches(1.5), Inches(5.4), Inches(0.42),
   size=13, bold=True, color=DARK_BLUE)
for i, (feat, ftype, rat) in enumerate(fe_items):
    y = Inches(2.0) + i * Inches(0.42)
    tb(s, f'{feat}', Inches(7.55), y, Inches(2.2), Inches(0.4), size=11, bold=True, color=DARK)
    tb(s, f'({ftype})', Inches(9.75), y, Inches(1.5), Inches(0.4), size=11, color=GREY)

# Result summary
box(s, Inches(7.4), Inches(5.35), Inches(5.7), Inches(1.75), fill=LT_ORANGE)
tb(s, 'Result', Inches(7.6), Inches(5.42), Inches(5.2), Inches(0.42),
   size=14, bold=True, color=ORANGE)
bullets(s, [
    'CV Macro F1: 0.5102 vs 0.5093 (+0.0009)',
    'Test Macro F1: 0.5093 vs 0.5100 (−0.0007)',
    '7/8 new features in RF Top 25 — but redundant',
    'Ceiling is data ambiguity, not missing features',
], Inches(7.5), Inches(5.88), Inches(5.5), Inches(1.15), size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusions
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header_bar(s, 'Conclusions & Future Work', '')
footer(s, 12)

# Two-column layout
# Left: conclusions
box(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.7), fill=LT_BLUE)
tb(s, 'Key Findings', Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
   size=17, bold=True, color=DARK_BLUE)
bullets(s, [
    'Random Forest is the best model — Test Macro F1 = 0.5100',
    'Non-linear methods (RF, XGBoost) outperform linear (LR, SVM) by +0.07',
    'SMOTE provides the only meaningful imbalance-handling gain (+0.029 on LR)',
    'Mixed→Good is the dominant error: 42.4 % of Mixed misclassified as Good',
    'Feature engineering (8 new features) adds no improvement — RF already captures count-level patterns',
    'Performance ceiling is inherent: Good and Mixed are genuinely similar in pre-release metadata space',
], Inches(0.55), Inches(2.1), Inches(6.0), Inches(4.8), size=15)

# Right: future work
box(s, Inches(6.9), Inches(1.4), Inches(6.1), Inches(5.7), fill=RGBColor(240,245,255))
tb(s, 'Future Directions', Inches(7.1), Inches(1.5), Inches(5.7), Inches(0.5),
   size=17, bold=True, color=DARK_BLUE)
bullets(s, [
    'NLP on game descriptions — semantic signals beyond metadata tags',
    'Developer reputation features — prior game performance as prior',
    'Post-launch enrichment — community tags, early reviews, patch frequency',
    'Ensemble stacking — combine RF and XGBoost probability outputs',
    'Finer threshold exploration — reduce genuinely borderline cases in the label',
    'SHAP interaction values — quantify feature pair synergies (e.g. 2D × era)',
], Inches(7.05), Inches(2.1), Inches(5.9), Inches(4.8), size=15)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = 'steam_games_presentation.pptx'
prs.save(OUT)
print(f'Saved: {OUT}')
